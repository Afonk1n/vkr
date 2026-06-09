# -*- coding: utf-8 -*-
"""Сборка презентации защиты ВКР Mustreview.

Запуск: python build_presentation.py
Выход: Mustreview_VKR.pptx рядом со скриптом.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).parent
ASSETS = ROOT / "_assets"
SCHEMAS = ROOT.parent / "Схемы" / "png"
OUTPUT = ROOT / "presentation_v1.pptx"

# Палитра — сине-белая, как у вуза
NAVY = RGBColor(0x1F, 0x3A, 0x8A)
NAVY_DARK = RGBColor(0x15, 0x29, 0x66)
NAVY_LIGHT = RGBColor(0x2E, 0x4D, 0xAA)
ACCENT_RED = RGBColor(0xC8, 0x10, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_50 = RGBColor(0xF4, 0xF6, 0xFB)
GRAY_BORDER = RGBColor(0xD8, 0xDE, 0xEB)
TEXT_DARK = RGBColor(0x1A, 0x1F, 0x36)
TEXT_MUTED = RGBColor(0x5A, 0x6B, 0x8C)

HEADER_FONT = "Cambria"
BODY_FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

THEME_TITLE = "Разработка информационной системы рецензирования музыкального творчества"

# ---------- Низкоуровневые помощники ----------


def add_rect(slide, x, y, w, h, fill=NAVY, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_rounded(slide, x, y, w, h, fill=WHITE, line=GRAY_BORDER):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.08
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    font=BODY_FONT,
    size=14,
    bold=False,
    italic=False,
    color=TEXT_DARK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_circle(slide, x, y, d, fill=NAVY):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


# ---------- Общие блоки для контентных слайдов ----------


def apply_content_chrome(slide, title, slide_number, total_slides):
    """Верхняя плашка с лого+заголовком, нижняя плашка с темой и номером."""
    # Верхний синий бар
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.95), fill=NAVY)
    # Логотип слева на баре
    logo = ASSETS / "mgtu_logo.png"
    if logo.exists():
        slide.shapes.add_picture(
            str(logo), Inches(0.32), Inches(0.08), height=Inches(0.78)
        )
    # Заголовок слайда
    add_text(
        slide,
        Inches(1.4),
        Inches(0.15),
        Inches(11.5),
        Inches(0.7),
        title,
        font=HEADER_FONT,
        size=26,
        bold=True,
        color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Нижняя полоса
    add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), fill=NAVY)
    add_text(
        slide,
        Inches(0.4),
        Inches(7.18),
        Inches(11),
        Inches(0.3),
        THEME_TITLE,
        font=BODY_FONT,
        size=10,
        italic=True,
        color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        Inches(11.8),
        Inches(7.18),
        Inches(1.3),
        Inches(0.3),
        f"{slide_number} / {total_slides}",
        font=BODY_FONT,
        size=10,
        color=WHITE,
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


# ---------- Слайды ----------


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Левая тёмно-синяя колонка
    add_rect(slide, 0, 0, Inches(4.4), SLIDE_H, fill=NAVY)
    add_rect(slide, Inches(4.4), 0, Inches(0.06), SLIDE_H, fill=ACCENT_RED)

    # Логотип
    logo = ASSETS / "mgtu_logo.png"
    if logo.exists():
        slide.shapes.add_picture(
            str(logo), Inches(1.55), Inches(0.6), height=Inches(2.6)
        )

    # Шапка вуза
    add_text(
        slide,
        Inches(0.35),
        Inches(3.4),
        Inches(3.7),
        Inches(2.6),
        [
            "МИНИСТЕРСТВО НАУКИ И ВЫСШЕГО",
            "ОБРАЗОВАНИЯ РОССИЙСКОЙ ФЕДЕРАЦИИ",
            "",
            "Магнитогорский государственный",
            "технический университет",
            "им. Г.И. Носова",
            "",
            "Институт энергетики и автоматизированных систем",
            "Кафедра вычислительной техники",
            "и программирования",
        ],
        font=BODY_FONT,
        size=10.5,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )

    # Декоративная плашка снизу левой колонки
    add_rect(slide, 0, Inches(6.7), Inches(4.4), Inches(0.8), fill=NAVY_DARK)
    add_text(
        slide,
        Inches(0.35),
        Inches(6.75),
        Inches(3.7),
        Inches(0.7),
        "Магнитогорск, 2026",
        font=BODY_FONT,
        size=12,
        italic=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Правая часть — заголовок работы
    add_text(
        slide,
        Inches(5.0),
        Inches(0.7),
        Inches(8.0),
        Inches(0.5),
        "ВЫПУСКНАЯ КВАЛИФИКАЦИОННАЯ РАБОТА",
        font=BODY_FONT,
        size=12,
        bold=True,
        color=NAVY,
    )
    add_text(
        slide,
        Inches(5.0),
        Inches(1.15),
        Inches(8.0),
        Inches(0.35),
        "бакалавра",
        font=BODY_FONT,
        size=11,
        italic=True,
        color=TEXT_MUTED,
    )

    # Акцентная плашка с темой
    add_rect(slide, Inches(5.0), Inches(1.8), Inches(0.12), Inches(2.6), fill=ACCENT_RED)
    add_text(
        slide,
        Inches(5.3),
        Inches(1.85),
        Inches(7.7),
        Inches(0.45),
        "Тема работы",
        font=BODY_FONT,
        size=11,
        bold=True,
        color=TEXT_MUTED,
    )
    add_text(
        slide,
        Inches(5.3),
        Inches(2.25),
        Inches(7.7),
        Inches(2.2),
        "Разработка информационной системы рецензирования музыкального творчества",
        font=HEADER_FONT,
        size=26,
        bold=True,
        color=NAVY_DARK,
    )

    # Блок автора/руководителя
    add_rounded(
        slide,
        Inches(5.0),
        Inches(4.8),
        Inches(8.0),
        Inches(1.9),
        fill=GRAY_50,
        line=GRAY_BORDER,
    )
    # Две колонки
    add_text(
        slide,
        Inches(5.3),
        Inches(4.95),
        Inches(3.7),
        Inches(0.3),
        "ВЫПОЛНИЛ",
        font=BODY_FONT,
        size=9,
        bold=True,
        color=TEXT_MUTED,
    )
    add_text(
        slide,
        Inches(5.3),
        Inches(5.25),
        Inches(3.7),
        Inches(0.4),
        "студент группы [укажите]",
        font=BODY_FONT,
        size=13,
        color=TEXT_DARK,
    )
    add_text(
        slide,
        Inches(5.3),
        Inches(5.65),
        Inches(3.7),
        Inches(0.5),
        "Афонкин М.А.",
        font=HEADER_FONT,
        size=18,
        bold=True,
        color=NAVY_DARK,
    )

    add_text(
        slide,
        Inches(9.2),
        Inches(4.95),
        Inches(3.7),
        Inches(0.3),
        "РУКОВОДИТЕЛЬ",
        font=BODY_FONT,
        size=9,
        bold=True,
        color=TEXT_MUTED,
    )
    add_text(
        slide,
        Inches(9.2),
        Inches(5.25),
        Inches(3.7),
        Inches(0.4),
        "[уч. степень, звание]",
        font=BODY_FONT,
        size=13,
        color=TEXT_DARK,
    )
    add_text(
        slide,
        Inches(9.2),
        Inches(5.65),
        Inches(3.7),
        Inches(0.5),
        "[Фамилия И.О.]",
        font=HEADER_FONT,
        size=18,
        bold=True,
        color=NAVY_DARK,
    )

    # Нижняя строка с темой проекта
    add_text(
        slide,
        Inches(5.0),
        Inches(6.95),
        Inches(8.0),
        Inches(0.4),
        "Mustreview · веб-сервис рецензий, оценок и геймификации",
        font=BODY_FONT,
        size=11,
        italic=True,
        color=TEXT_MUTED,
        align=PP_ALIGN.RIGHT,
    )


def slide_goal(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_content_chrome(slide, "Цель, объект и предмет исследования", 2, total)

    cards = [
        (
            "🎯",
            "ЦЕЛЬ",
            "Разработать веб-приложение для публикации, оценки и анализа рецензий на музыкальные релизы с системой геймификации и модерацией пользовательского контента.",
        ),
        (
            "🔍",
            "ОБЪЕКТ",
            "Процессы рецензирования музыкального творчества в онлайн-сообществах: создание, оценка, обсуждение и распространение пользовательских рецензий.",
        ),
        (
            "📌",
            "ПРЕДМЕТ",
            "Программные и алгоритмические средства реализации информационной системы рецензирования: модель данных, REST API, конвейер модерации и алгоритмы оценки.",
        ),
    ]

    card_w = Inches(4.0)
    gap = Inches(0.2)
    start_x = (SLIDE_W - card_w * 3 - gap * 2) // 2
    top = Inches(1.55)
    card_h = Inches(5.2)

    for i, (emoji, title, body) in enumerate(cards):
        x = start_x + (card_w + gap) * i
        add_rounded(slide, x, top, card_w, card_h, fill=WHITE, line=GRAY_BORDER)
        # цветная шапка
        add_rect(slide, x, top, card_w, Inches(0.18), fill=NAVY)
        # иконка
        add_circle(slide, x + Inches(0.4), top + Inches(0.5), Inches(0.9), fill=NAVY)
        add_text(
            slide,
            x + Inches(0.4),
            top + Inches(0.5),
            Inches(0.9),
            Inches(0.9),
            emoji,
            size=32,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            x + Inches(1.5),
            top + Inches(0.6),
            card_w - Inches(1.7),
            Inches(0.6),
            title,
            font=HEADER_FONT,
            size=22,
            bold=True,
            color=NAVY_DARK,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            x + Inches(0.4),
            top + Inches(1.7),
            card_w - Inches(0.8),
            Inches(3.3),
            body,
            font=BODY_FONT,
            size=14,
            color=TEXT_DARK,
        )


def slide_tasks(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_content_chrome(slide, "Задачи исследования", 3, total)

    tasks = [
        "Провести анализ существующих сервисов музыкального рецензирования (RateYourMusic, Last.fm, Genius) и сформулировать функциональные требования.",
        "Спроектировать архитектуру системы и логическую модель базы данных, обеспечивающую хранение пользовательского контента и связей сообщества.",
        "Спроектировать алгоритмы расчёта итоговой оценки релиза и накопления опыта пользователя в системе геймификации.",
        "Реализовать клиент-серверное приложение на стеке React + Go (Gin) + PostgreSQL с REST API.",
        "Реализовать конвейер модерации пользовательских рецензий и механизм управления правами доступа.",
        "Подготовить контейнеризацию (Docker Compose) и CI/CD-пайплайн с автоматическими проверками и публикацией образов.",
    ]

    col_w = Inches(5.85)
    gap = Inches(0.3)
    top = Inches(1.4)
    row_h = Inches(1.65)

    for i, text in enumerate(tasks):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + (col_w + gap) * col
        y = top + row_h * row
        # карточка задачи
        add_rounded(slide, x, y, col_w, row_h - Inches(0.18), fill=WHITE, line=GRAY_BORDER)
        # номер в круге
        add_circle(slide, x + Inches(0.25), y + Inches(0.32), Inches(0.85), fill=NAVY)
        add_text(
            slide,
            x + Inches(0.25),
            y + Inches(0.32),
            Inches(0.85),
            Inches(0.85),
            str(i + 1),
            font=HEADER_FONT,
            size=26,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            x + Inches(1.25),
            y + Inches(0.22),
            col_w - Inches(1.45),
            row_h - Inches(0.4),
            text,
            font=BODY_FONT,
            size=13.5,
            color=TEXT_DARK,
            anchor=MSO_ANCHOR.MIDDLE,
        )


def slide_analytics(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_content_chrome(slide, "Результаты аналитического исследования", 4, total)

    # Левая часть: сравнительная таблица аналогов
    add_text(
        slide,
        Inches(0.6),
        Inches(1.25),
        Inches(7.6),
        Inches(0.4),
        "Анализ существующих сервисов",
        font=HEADER_FONT,
        size=18,
        bold=True,
        color=NAVY_DARK,
    )

    # Шапка таблицы
    table_x = Inches(0.6)
    table_y = Inches(1.85)
    col_widths = [Inches(2.4), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3)]
    row_h = Inches(0.55)

    headers = ["Сервис", "Рецензии", "Соц-граф", "Геймиф.", "RU-интерф."]
    rows = [
        ["RateYourMusic", "✓", "ограничен", "—", "—"],
        ["Last.fm", "ограниченно", "✓", "—", "ограничен"],
        ["Genius", "только текст", "—", "—", "—"],
        ["Spotify", "—", "✓", "—", "✓"],
        ["Mustreview", "✓", "✓", "✓", "✓"],
    ]

    # Header row
    x = table_x
    for i, h in enumerate(headers):
        add_rect(slide, x, table_y, col_widths[i], row_h, fill=NAVY)
        add_text(
            slide,
            x,
            table_y,
            col_widths[i],
            row_h,
            h,
            font=BODY_FONT,
            size=11,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        x += col_widths[i]

    for r_idx, row in enumerate(rows):
        is_us = row[0] == "Mustreview"
        x = table_x
        y = table_y + row_h * (r_idx + 1)
        fill = ACCENT_RED if is_us else (GRAY_50 if r_idx % 2 == 0 else WHITE)
        text_color = WHITE if is_us else TEXT_DARK
        for c_idx, cell in enumerate(row):
            add_rect(slide, x, y, col_widths[c_idx], row_h, fill=fill, line=GRAY_BORDER)
            add_text(
                slide,
                x,
                y,
                col_widths[c_idx],
                row_h,
                cell,
                font=BODY_FONT,
                size=11,
                bold=is_us or c_idx == 0,
                color=text_color,
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            x += col_widths[c_idx]

    # Правая часть: обоснование стека
    add_text(
        slide,
        Inches(8.7),
        Inches(1.25),
        Inches(4.2),
        Inches(0.4),
        "Обоснование стека",
        font=HEADER_FONT,
        size=18,
        bold=True,
        color=NAVY_DARK,
    )

    stack_items = [
        ("React 18", "интерактивный UI с большим количеством состояний"),
        ("Go + Gin", "компактное и быстрое REST API, единый бинарь"),
        ("PostgreSQL", "связанные данные, транзакции, JSON-поля"),
        ("Docker Compose", "воспроизводимый запуск и единый артефакт"),
        ("GitHub Actions", "автоматические проверки и публикация образов"),
    ]
    top = Inches(1.85)
    for i, (name, desc) in enumerate(stack_items):
        y = top + Inches(0.85) * i
        add_rounded(
            slide,
            Inches(8.7),
            y,
            Inches(4.2),
            Inches(0.75),
            fill=GRAY_50,
            line=GRAY_BORDER,
        )
        add_rect(slide, Inches(8.7), y, Inches(0.12), Inches(0.75), fill=NAVY)
        add_text(
            slide,
            Inches(8.95),
            y + Inches(0.08),
            Inches(3.9),
            Inches(0.35),
            name,
            font=BODY_FONT,
            size=13,
            bold=True,
            color=NAVY_DARK,
        )
        add_text(
            slide,
            Inches(8.95),
            y + Inches(0.4),
            Inches(3.9),
            Inches(0.35),
            desc,
            font=BODY_FONT,
            size=11,
            color=TEXT_MUTED,
        )


def slide_er(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_content_chrome(slide, "Логическая модель базы данных", 5, total)

    er = SCHEMAS / "02_er_model.png"
    if er.exists():
        slide.shapes.add_picture(
            str(er),
            Inches(0.6),
            Inches(1.3),
            width=Inches(9.0),
        )
    else:
        add_rounded(
            slide,
            Inches(0.6),
            Inches(1.3),
            Inches(9.0),
            Inches(5.4),
            fill=GRAY_50,
            line=GRAY_BORDER,
        )
        add_text(
            slide,
            Inches(0.6),
            Inches(3.7),
            Inches(9.0),
            Inches(0.6),
            "[ER-диаграмма: schema/02_er_model.png]",
            font=BODY_FONT,
            size=14,
            italic=True,
            color=TEXT_MUTED,
            align=PP_ALIGN.CENTER,
        )

    # Боковая колонка — описание сущностей
    add_text(
        slide,
        Inches(9.9),
        Inches(1.3),
        Inches(3.0),
        Inches(0.4),
        "Ключевые сущности",
        font=HEADER_FONT,
        size=16,
        bold=True,
        color=NAVY_DARK,
    )

    entities = [
        ("users", "учётные записи, профиль, роли"),
        ("albums / tracks", "каталог музыкальных объектов"),
        ("genres", "жанровая классификация"),
        ("reviews", "оценки и тексты, статус модерации"),
        ("*_likes", "лайки альбомов, треков, рецензий"),
        ("user_follows", "социальный граф подписок"),
    ]
    top = Inches(1.85)
    for i, (name, desc) in enumerate(entities):
        y = top + Inches(0.78) * i
        add_rect(slide, Inches(9.9), y, Inches(0.08), Inches(0.7), fill=ACCENT_RED)
        add_text(
            slide,
            Inches(10.05),
            y,
            Inches(2.9),
            Inches(0.35),
            name,
            font=BODY_FONT,
            size=12,
            bold=True,
            color=NAVY_DARK,
        )
        add_text(
            slide,
            Inches(10.05),
            y + Inches(0.32),
            Inches(2.9),
            Inches(0.4),
            desc,
            font=BODY_FONT,
            size=10.5,
            color=TEXT_MUTED,
        )


def slide_algorithms(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_content_chrome(slide, "Ключевые алгоритмы", 6, total)

    blocks = [
        {
            "title": "Расчёт итоговой оценки",
            "kpi": "1–90",
            "kpi_label": "шкала итога",
            "body": (
                "4 параметра оценки (рифмы, структура, реализация, харизма) — каждый 1–10. "
                "Атмосфера переводится в множитель. Итог = (сумма параметров) × множитель, "
                "приведённый к шкале 1–90. Формула скрыта от пользователя — показывается крупный итог."
            ),
        },
        {
            "title": "Расчёт уровня профиля",
            "kpi": "5",
            "kpi_label": "уровней званий",
            "body": (
                "Опыт начисляется за активность, а НЕ за высоту оценок: +320 за рецензию, "
                "+55 за лайк своей рецензии, +12 за поставленный лайк, +240 за авторский лайк. "
                "Уровни: бронза → серебро → золото → изумруд → платина (0 / 2 500 / 8 000 / 18 000 / 36 000)."
            ),
        },
        {
            "title": "Конвейер модерации",
            "kpi": "3",
            "kpi_label": "статуса жизни",
            "body": (
                "Новая рецензия → pending; админ принимает или отклоняет. В ленту попадают только "
                "approved. После approve пересчитывается средняя оценка альбома и опыт автора. "
                "Доступ к модерации защищён AdminMiddleware и подписанным сессионным токеном."
            ),
        },
    ]

    block_w = Inches(4.0)
    gap = Inches(0.2)
    start_x = (SLIDE_W - block_w * 3 - gap * 2) // 2
    top = Inches(1.4)
    block_h = Inches(5.4)

    for i, b in enumerate(blocks):
        x = start_x + (block_w + gap) * i
        add_rounded(slide, x, top, block_w, block_h, fill=WHITE, line=GRAY_BORDER)
        # цветная шапка блока
        add_rect(slide, x, top, block_w, Inches(0.7), fill=NAVY)
        add_text(
            slide,
            x + Inches(0.3),
            top,
            block_w - Inches(0.6),
            Inches(0.7),
            b["title"],
            font=HEADER_FONT,
            size=16,
            bold=True,
            color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # KPI
        add_text(
            slide,
            x + Inches(0.3),
            top + Inches(0.95),
            block_w - Inches(0.6),
            Inches(1.2),
            b["kpi"],
            font=HEADER_FONT,
            size=72,
            bold=True,
            color=ACCENT_RED,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            x + Inches(0.3),
            top + Inches(2.15),
            block_w - Inches(0.6),
            Inches(0.4),
            b["kpi_label"],
            font=BODY_FONT,
            size=11,
            italic=True,
            color=TEXT_MUTED,
            align=PP_ALIGN.CENTER,
        )
        # разделитель
        add_rect(
            slide,
            x + Inches(1.0),
            top + Inches(2.7),
            block_w - Inches(2.0),
            Emu(15875),
            fill=GRAY_BORDER,
        )
        # описание
        add_text(
            slide,
            x + Inches(0.3),
            top + Inches(2.9),
            block_w - Inches(0.6),
            Inches(2.4),
            b["body"],
            font=BODY_FONT,
            size=12,
            color=TEXT_DARK,
        )


def slide_results(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_content_chrome(slide, "Полученные результаты", 7, total)

    # Stat-callouts вверху
    stats = [
        ("35+", "REST-эндпоинтов"),
        ("8", "таблиц БД"),
        ("4", "слайда CI-пайплайна"),
        ("1", "команда запуска"),
    ]
    stat_w = Inches(2.95)
    gap = Inches(0.2)
    start_x = (SLIDE_W - stat_w * 4 - gap * 3) // 2
    top = Inches(1.3)
    stat_h = Inches(1.5)
    for i, (kpi, label) in enumerate(stats):
        x = start_x + (stat_w + gap) * i
        add_rounded(slide, x, top, stat_w, stat_h, fill=NAVY, line=None)
        add_text(
            slide,
            x,
            top + Inches(0.1),
            stat_w,
            Inches(0.9),
            kpi,
            font=HEADER_FONT,
            size=48,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            x,
            top + Inches(1.0),
            stat_w,
            Inches(0.4),
            label,
            font=BODY_FONT,
            size=12,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Левая колонка — реализовано
    list_top = Inches(3.05)
    add_text(
        slide,
        Inches(0.6),
        list_top,
        Inches(6.0),
        Inches(0.4),
        "Реализовано",
        font=HEADER_FONT,
        size=18,
        bold=True,
        color=NAVY_DARK,
    )

    items = [
        "Авторизация по подписанным сессионным токенам, роли admin/user",
        "Профиль пользователя с уровнями, бейджами, предпочтениями и социальным графом",
        "Лента рецензий с фильтром «Все / Подписки» и админ-панель модерации",
        "Лайки альбомов / треков / рецензий с оптимистичным UI",
        "Идемпотентный сидер демо-данных для воспроизводимой демонстрации",
        "GitHub Actions: go vet / test / build + smoke compose + публикация образов в GHCR",
    ]
    y = list_top + Inches(0.55)
    for txt in items:
        add_text(
            slide,
            Inches(0.7),
            y,
            Inches(0.3),
            Inches(0.35),
            "▸",
            font=BODY_FONT,
            size=13,
            bold=True,
            color=ACCENT_RED,
        )
        add_text(
            slide,
            Inches(1.0),
            y,
            Inches(5.5),
            Inches(0.55),
            txt,
            font=BODY_FONT,
            size=12.5,
            color=TEXT_DARK,
        )
        y += Inches(0.5)

    # Правая колонка — схема архитектуры или CI
    arch = SCHEMAS / "01_architecture.png"
    add_rounded(
        slide,
        Inches(6.9),
        list_top,
        Inches(5.9),
        Inches(3.6),
        fill=GRAY_50,
        line=GRAY_BORDER,
    )
    add_text(
        slide,
        Inches(7.1),
        list_top + Inches(0.12),
        Inches(5.5),
        Inches(0.35),
        "Архитектура системы",
        font=BODY_FONT,
        size=12,
        bold=True,
        color=NAVY_DARK,
    )
    if arch.exists():
        slide.shapes.add_picture(
            str(arch),
            Inches(7.05),
            list_top + Inches(0.5),
            width=Inches(5.6),
        )
    else:
        add_text(
            slide,
            Inches(6.9),
            list_top + Inches(1.4),
            Inches(5.9),
            Inches(0.5),
            "[Схема: Пояснительная записка/Схемы/png/01_architecture.png]",
            font=BODY_FONT,
            size=11,
            italic=True,
            color=TEXT_MUTED,
            align=PP_ALIGN.CENTER,
        )


def slide_thanks(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Сплошной тёмно-синий фон
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY_DARK)
    # Красный акцент
    add_rect(slide, 0, Inches(3.55), SLIDE_W, Inches(0.06), fill=ACCENT_RED)

    # Лого
    logo = ASSETS / "mgtu_logo.png"
    if logo.exists():
        slide.shapes.add_picture(
            str(logo), Inches(6.27), Inches(0.8), height=Inches(2.2)
        )

    # Основной текст
    add_text(
        slide,
        Inches(0.5),
        Inches(3.85),
        Inches(12.3),
        Inches(1.4),
        "СПАСИБО ЗА ВНИМАНИЕ!",
        font=HEADER_FONT,
        size=64,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        Inches(0.5),
        Inches(5.4),
        Inches(12.3),
        Inches(0.5),
        "Готов ответить на ваши вопросы",
        font=BODY_FONT,
        size=20,
        italic=True,
        color=GRAY_50,
        align=PP_ALIGN.CENTER,
    )

    # Подвал
    add_text(
        slide,
        Inches(0.5),
        Inches(6.7),
        Inches(12.3),
        Inches(0.4),
        "Афонкин М.А.  ·  МГТУ им. Г.И. Носова  ·  Магнитогорск, 2026",
        font=BODY_FONT,
        size=12,
        color=GRAY_50,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        Inches(0.5),
        Inches(7.1),
        Inches(12.3),
        Inches(0.3),
        "Mustreview · веб-сервис рецензий, оценок и геймификации",
        font=BODY_FONT,
        size=11,
        italic=True,
        color=GRAY_BORDER,
        align=PP_ALIGN.CENTER,
    )


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,
        slide_goal,
        slide_tasks,
        slide_analytics,
        slide_er,
        slide_algorithms,
        slide_results,
        slide_thanks,
    ]
    total = len(builders)
    # Первый слайд без chrome'а, остальные — со счётчиком
    slide_title(prs)
    slide_goal(prs, total)
    slide_tasks(prs, total)
    slide_analytics(prs, total)
    slide_er(prs, total)
    slide_algorithms(prs, total)
    slide_results(prs, total)
    slide_thanks(prs, total)

    prs.save(str(OUTPUT))
    print(f"saved: {OUTPUT}")


if __name__ == "__main__":
    main()
