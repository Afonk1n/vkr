# -*- coding: utf-8 -*-
"""Сборка ФИНАЛЬНОЙ презентации защиты ВКР (синергия v1 + v2).

Запуск: python build_v3.py
Выход: presentation_v3.pptx рядом со скриптом.

Учтённое ТЗ:
- DevOps-слайд обязателен (перерисован нативно, читаемо);
- логотип полный, не обрезанный;
- цель/объект/предмет — формат с эмодзи, текст крупнее;
- задачи — текст крупнее;
- логическая модель — оформление как в v1 (схема + сайдбар), но ER перерисована
  нативно (ровные связи) и текст крупнее;
- слайд расчёта оценки — в стиле v2 (формула + критерии), диапазон 6-90;
- результаты — в стиле v1 (больше инфо, минимум скринов).
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).parent
ASSETS = ROOT / "_assets"
OUTPUT = ROOT / "presentation_v3.pptx"

# Палитра — сине-белая, как у вуза
NAVY = RGBColor(0x1F, 0x3A, 0x8A)
NAVY_DARK = RGBColor(0x15, 0x29, 0x66)
NAVY_LIGHT = RGBColor(0x2E, 0x4D, 0xAA)
ACCENT_RED = RGBColor(0xC8, 0x10, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_50 = RGBColor(0xF4, 0xF6, 0xFB)
GRAY_BORDER = RGBColor(0xD8, 0xDE, 0xEB)
TEXT_DARK = RGBColor(0x1A, 0x1F, 0x36)
TEXT_MUTED = RGBColor(0x5A, 0x6B, 0x8C)
LINK_GRAY = RGBColor(0x9A, 0xA6, 0xBF)

TITLE_FONT = "Calibri"   # рубленый, современный (как в v2)
BODY_FONT = "Calibri"
SERIF_FONT = "Cambria"   # для крупных акцентов (формула)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# --- Реальные данные титула (сверено с docx) ---
THEME = "Проектирование и разработка WEB-приложения для рецензирования музыкального творчества"
THEME_SHORT = "WEB-приложение для рецензирования музыкального творчества"
AUTHOR_FULL = "Афонькин Максим Артемович"
AUTHOR_SHORT = "Афонькин М. А."
GROUP = "АВб-22-2"
SUPERVISOR = "Егорова Л. Г."
SUPERVISOR_RANK = "Доцент каф. ВТ и П, к.т.н."
CITY_YEAR = "Магнитогорск, 2026"


# ---------- Низкоуровневые помощники ----------


def add_rect(slide, x, y, w, h, fill=NAVY, line=None, line_w=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_rounded(slide, x, y, w, h, fill=WHITE, line=GRAY_BORDER, radius=0.08, line_w=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_circle(slide, x, y, d, fill=NAVY):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, font=BODY_FONT, size=14, bold=False,
             italic=False, color=TEXT_DARK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_connector(slide, x1, y1, x2, y2, color=LINK_GRAY, width=1.5, elbow=False):
    ctype = MSO_CONNECTOR.ELBOW if elbow else MSO_CONNECTOR.STRAIGHT
    cn = slide.shapes.add_connector(ctype, x1, y1, x2, y2)
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    cn.shadow.inherit = False
    return cn


LOGO_RATIO = 0.56  # w / h


def logo(slide, x, y, h, badge=False, pad=Inches(0.13)):
    """Лого МГТУ. На тёмном фоне badge=True подкладывает белую плашку для контраста."""
    w = Emu(int(h * LOGO_RATIO))
    if badge:
        add_rounded(slide, x - pad, y - pad, w + pad * 2, h + pad * 2,
                    fill=WHITE, line=None, radius=0.12)
    p = ASSETS / "mgtu_logo_full.png"
    if not p.exists():
        p = ASSETS / "mgtu_logo.png"
    if p.exists():
        slide.shapes.add_picture(str(p), x, y, height=h)


# ---------- Chrome контентных слайдов ----------


def chrome(slide, title, number, total):
    # Верхний бар (повыше, чтобы логотип влезал целиком)
    bar_h = Inches(1.02)
    add_rect(slide, 0, 0, SLIDE_W, bar_h, fill=NAVY)
    add_rect(slide, 0, bar_h, SLIDE_W, Inches(0.05), fill=ACCENT_RED)
    # Логотип в белом значке (контраст на синем баре)
    logo(slide, Inches(0.4), Inches(0.16), Inches(0.7), badge=True, pad=Inches(0.1))
    add_text(slide, Inches(1.35), 0, Inches(11.4), bar_h, title,
             font=TITLE_FONT, size=30, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    # Нижняя полоса
    add_rect(slide, 0, Inches(7.12), SLIDE_W, Inches(0.38), fill=NAVY)
    add_text(slide, Inches(0.4), Inches(7.15), Inches(11), Inches(0.32),
             THEME_SHORT, font=BODY_FONT, size=10, italic=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(11.7), Inches(7.15), Inches(1.4), Inches(0.32),
             f"{number} / {total}", font=BODY_FONT, size=10, color=WHITE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------- Слайд 1: Титул ----------


def slide_title(prs):
    s = blank(prs)
    add_rect(s, 0, 0, Inches(4.5), SLIDE_H, fill=NAVY)
    add_rect(s, Inches(4.5), 0, Inches(0.07), SLIDE_H, fill=ACCENT_RED)

    logo(s, Inches(1.6), Inches(0.6), Inches(2.5), badge=True, pad=Inches(0.18))

    add_text(s, Inches(0.35), Inches(3.45), Inches(3.8), Inches(2.7),
             ["МИНИСТЕРСТВО НАУКИ И ВЫСШЕГО",
              "ОБРАЗОВАНИЯ РОССИЙСКОЙ ФЕДЕРАЦИИ", "",
              "Магнитогорский государственный",
              "технический университет им. Г.И. Носова", "",
              "Институт энергетики и",
              "автоматизированных систем",
              "Кафедра вычислительной техники",
              "и программирования"],
             font=BODY_FONT, size=11, color=WHITE, align=PP_ALIGN.CENTER,
             line_spacing=1.1)

    add_rect(s, 0, Inches(6.72), Inches(4.5), Inches(0.78), fill=NAVY_DARK)
    add_text(s, Inches(0.35), Inches(6.72), Inches(3.8), Inches(0.78),
             CITY_YEAR, font=BODY_FONT, size=13, italic=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Правая часть
    add_text(s, Inches(5.05), Inches(0.65), Inches(8.0), Inches(0.45),
             "ВЫПУСКНАЯ КВАЛИФИКАЦИОННАЯ РАБОТА", font=BODY_FONT, size=13,
             bold=True, color=NAVY)
    add_text(s, Inches(5.05), Inches(1.1), Inches(8.0), Inches(0.35),
             "бакалаврская работа", font=BODY_FONT, size=11, italic=True,
             color=TEXT_MUTED)

    add_rect(s, Inches(5.05), Inches(1.75), Inches(0.13), Inches(2.7), fill=ACCENT_RED)
    add_text(s, Inches(5.35), Inches(1.78), Inches(7.6), Inches(0.4),
             "Тема работы", font=BODY_FONT, size=12, bold=True, color=TEXT_MUTED)
    add_text(s, Inches(5.35), Inches(2.2), Inches(7.65), Inches(2.3),
             THEME, font=TITLE_FONT, size=27, bold=True, color=NAVY_DARK,
             line_spacing=1.05)

    # Блок автора/руководителя
    add_rounded(s, Inches(5.05), Inches(4.95), Inches(7.95), Inches(1.95),
                fill=GRAY_50, line=GRAY_BORDER)
    add_text(s, Inches(5.35), Inches(5.12), Inches(3.6), Inches(0.3),
             "ВЫПОЛНИЛ", font=BODY_FONT, size=10, bold=True, color=TEXT_MUTED)
    add_text(s, Inches(5.35), Inches(5.42), Inches(3.6), Inches(0.4),
             f"студент группы {GROUP}", font=BODY_FONT, size=13, color=TEXT_DARK)
    add_text(s, Inches(5.35), Inches(5.85), Inches(3.6), Inches(0.5),
             AUTHOR_SHORT, font=TITLE_FONT, size=19, bold=True, color=NAVY_DARK)

    add_text(s, Inches(9.25), Inches(5.12), Inches(3.6), Inches(0.3),
             "РУКОВОДИТЕЛЬ", font=BODY_FONT, size=10, bold=True, color=TEXT_MUTED)
    add_text(s, Inches(9.25), Inches(5.42), Inches(3.6), Inches(0.4),
             SUPERVISOR_RANK, font=BODY_FONT, size=13, color=TEXT_DARK)
    add_text(s, Inches(9.25), Inches(5.85), Inches(3.6), Inches(0.5),
             SUPERVISOR, font=TITLE_FONT, size=18, bold=True, color=NAVY_DARK)


# ---------- Слайд 2: Цель, объект, предмет ----------


def slide_goal(prs, total):
    s = blank(prs)
    chrome(s, "Цель, объект и предмет исследования", 2, total)

    cards = [
        ("🎯", "ЦЕЛЬ",
         "Повышение объективности, информативности и доступности рецензирования "
         "русскоязычных музыкальных произведений молодёжной аудиторией за счёт "
         "разработки специализированной информационной системы "
         "многокритериального оценивания."),
        ("🔍", "ОБЪЕКТ",
         "Процесс коллективного рецензирования музыкального творчества "
         "русскоязычной молодёжью на веб-платформах."),
        ("📌", "ПРЕДМЕТ",
         "Информационная система рецензирования современной музыки с экспертной "
         "системой оценки и инструментами модерации пользовательского контента."),
    ]
    card_w = Inches(4.0)
    gap = Inches(0.22)
    start_x = (SLIDE_W - card_w * 3 - gap * 2) // 2
    top = Inches(1.5)
    card_h = Inches(5.35)

    for i, (emoji, title, body) in enumerate(cards):
        x = start_x + (card_w + gap) * i
        add_rounded(s, x, top, card_w, card_h, fill=WHITE, line=GRAY_BORDER)
        add_rect(s, x, top, card_w, Inches(0.2), fill=NAVY)
        # иконка
        add_circle(s, x + Inches(0.35), top + Inches(0.55), Inches(1.1), fill=NAVY)
        add_text(s, x + Inches(0.35), top + Inches(0.55), Inches(1.1), Inches(1.1),
                 emoji, size=40, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(1.6), top + Inches(0.6), card_w - Inches(1.8),
                 Inches(1.0), title, font=TITLE_FONT, size=27, bold=True,
                 color=NAVY_DARK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.4), top + Inches(1.9), card_w - Inches(0.8),
                 Inches(3.3), body, font=BODY_FONT, size=16.5, color=TEXT_DARK,
                 line_spacing=1.12)


# ---------- Слайд 3: Задачи ----------


def slide_tasks(prs, total):
    s = blank(prs)
    chrome(s, "Задачи исследования", 3, total)

    add_text(s, Inches(0.55), Inches(1.3), Inches(12.2), Inches(0.45),
             "Для достижения поставленной цели решаются следующие задачи:",
             font=BODY_FONT, size=16, italic=True, color=TEXT_MUTED)

    tasks = [
        ("Анализ предметной области",
         "Анализ ниши русскоязычного рецензирования современной музыки "
         "и проблем существующих платформ."),
        ("Проектирование системы",
         "Проектирование архитектуры системы с экспертной оценкой "
         "(рифмы/образы, структура/ритмика, реализация стиля, индивидуальность/харизма, "
         "атмосфера/вайб), инструментами модерации и фильтрами жанров."),
        ("Реализация прототипа",
         "Реализация прототипа, интеграция каталогов музыкальных произведений, "
         "тестирование модерации и пользовательского интерфейса."),
    ]
    x = Inches(0.55)
    w = Inches(12.2)
    top = Inches(1.95)
    row_h = Inches(1.5)
    gap = Inches(0.25)
    for i, (head, body) in enumerate(tasks):
        y = top + (row_h + gap) * i
        add_rounded(s, x, y, w, row_h, fill=WHITE, line=GRAY_BORDER)
        add_circle(s, x + Inches(0.35), y + Inches(0.45), Inches(0.95), fill=NAVY)
        add_text(s, x + Inches(0.35), y + Inches(0.45), Inches(0.95), Inches(0.95),
                 str(i + 1), font=TITLE_FONT, size=34, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(1.6), y + Inches(0.16), w - Inches(1.9), Inches(0.45),
                 head, font=TITLE_FONT, size=20, bold=True, color=NAVY_DARK)
        add_text(s, x + Inches(1.6), y + Inches(0.64), w - Inches(1.9), Inches(0.8),
                 body, font=BODY_FONT, size=16.5, color=TEXT_DARK, line_spacing=1.1)


# ---------- Слайд 4: Аналитика ----------


def slide_analytics(prs, total):
    s = blank(prs)
    chrome(s, "Результаты аналитического исследования", 4, total)

    add_text(s, Inches(0.55), Inches(1.3), Inches(7.6), Inches(0.45),
             "Сравнение существующих сервисов", font=TITLE_FONT, size=19,
             bold=True, color=NAVY_DARK)

    table_x = Inches(0.55)
    table_y = Inches(1.95)
    col_w = [Inches(2.5), Inches(1.35), Inches(1.35), Inches(1.35), Inches(1.35)]
    row_h = Inches(0.62)
    headers = ["Сервис", "Много-\nкритер.", "Модер.\nконтента", "Гейми-\nфикация", "RU-\nинтерфейс"]
    rows = [
        ["RateYourMusic", "—", "±", "—", "—"],
        ["AllMusic", "—", "±", "—", "—"],
        ["Last.fm", "—", "—", "—", "±"],
        ["VK / Я.Музыка", "—", "—", "—", "+"],
        ["Mustreview", "+", "+", "+", "+"],
    ]
    # header
    x = table_x
    for i, h in enumerate(headers):
        add_rect(s, x, table_y, col_w[i], row_h, fill=NAVY)
        add_text(s, x, table_y, col_w[i], row_h, h.split("\n"), font=BODY_FONT,
                 size=12.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[i]
    for r, row in enumerate(rows):
        is_us = row[0] == "Mustreview"
        y = table_y + row_h * (r + 1)
        x = table_x
        fill = ACCENT_RED if is_us else (GRAY_50 if r % 2 == 0 else WHITE)
        tcol = WHITE if is_us else TEXT_DARK
        for c, cell in enumerate(row):
            add_rect(s, x, y, col_w[c], row_h, fill=fill, line=GRAY_BORDER, line_w=0.5)
            add_text(s, x, y, col_w[c], row_h, cell, font=BODY_FONT,
                     size=14.5, bold=is_us or c == 0, color=tcol,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            x += col_w[c]

    # Правая часть: ключевые выводы
    rx = Inches(8.7)
    add_text(s, rx, Inches(1.3), Inches(4.2), Inches(0.45), "Ключевые выводы",
             font=TITLE_FONT, size=19, bold=True, color=NAVY_DARK)
    findings = [
        ("Дефицит платформ", "русскоязычный сегмент фрагментирован, нет среды для системного рецензирования"),
        ("Технологический стек", "Go + Gin, React (SPA), PostgreSQL, Docker + GitHub Actions"),
        ("Главное отличие", "5-критериальная оценка + очередь модерации + отметка от исполнителя"),
    ]
    ty = Inches(1.95)
    for name, desc in findings:
        add_rounded(s, rx, ty, Inches(4.2), Inches(1.2), fill=GRAY_50, line=GRAY_BORDER)
        add_rect(s, rx, ty, Inches(0.13), Inches(1.2), fill=NAVY)
        add_text(s, rx + Inches(0.3), ty + Inches(0.13), Inches(3.8), Inches(0.4),
                 name, font=BODY_FONT, size=15.5, bold=True, color=NAVY_DARK)
        add_text(s, rx + Inches(0.3), ty + Inches(0.52), Inches(3.75), Inches(0.6),
                 desc, font=BODY_FONT, size=13.5, color=TEXT_MUTED, line_spacing=1.1)
        ty += Inches(1.4)


# ---------- Слайд 5: Логическая модель БД (нативная ER) ----------


def slide_db(prs, total):
    s = blank(prs)
    chrome(s, "Физическая модель базы данных", 5, total)

    add_text(s, Inches(0.55), Inches(1.25), Inches(12.2), Inches(0.4),
             "PostgreSQL · 5 таблиц · связи реализованы внешними ключами (FK)",
             font=BODY_FONT, size=16, italic=True, color=TEXT_MUTED)

    def card(x, y, w, h, emoji, name, desc, fields, accent=NAVY):
        add_rounded(s, x, y, w, h, fill=WHITE, line=GRAY_BORDER, line_w=1.0)
        # цветная шапка
        add_rect(s, x, y, w, Inches(0.74), fill=accent)
        add_circle(s, x + Inches(0.18), y + Inches(0.13), Inches(0.48), fill=WHITE)
        add_text(s, x + Inches(0.18), y + Inches(0.13), Inches(0.48), Inches(0.48),
                 emoji, size=20, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.78), y, w - Inches(0.9), Inches(0.74), name,
                 font="Consolas", size=18, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        # описание
        add_text(s, x + Inches(0.22), y + Inches(0.85), w - Inches(0.44), Inches(0.55),
                 desc, font=BODY_FONT, size=14, bold=True, color=NAVY_DARK,
                 line_spacing=1.05)
        # поля
        add_text(s, x + Inches(0.22), y + Inches(1.42), w - Inches(0.44), h - Inches(1.5),
                 fields, font="Consolas", size=12.5, color=TEXT_MUTED, line_spacing=1.2)

    # верхний ряд: 3 карточки
    w1 = Inches(4.05)
    gap = Inches(0.235)
    x0 = (SLIDE_W - w1 * 3 - gap * 2) // 2
    y1 = Inches(1.78)
    h1 = Inches(2.35)
    card(x0, y1, w1, h1, "👤", "users", "Пользователи и роли",
         "id · username · email\npassword · avatar_path\nbio · is_admin · created_at")
    card(x0 + (w1 + gap), y1, w1, h1, "💿", "albums", "Каталог альбомов",
         "id · genre_id (FK)\ntitle · artist · cover_url\ndescription · created_at")
    card(x0 + (w1 + gap) * 2, y1, w1, h1, "🎵", "tracks", "Треки альбомов",
         "id · album_id (FK)\ntitle · duration\ntrack_number · created_at")

    # нижний ряд: genres (узкая) + reviews (широкая, акцентная)
    y2 = Inches(4.35)
    h2 = Inches(2.55)
    wg = Inches(4.05)
    wr = Inches(8.29)
    card(x0, y2, wg, h2, "🏷️", "genres", "Жанровая классификация",
         "id · name · description")
    card(x0 + (wg + gap), y2, wr, h2, "⭐", "reviews",
         "Центральная сущность — рецензии и оценки",
         "id · user_id (FK) · album_id (FK) · track_id (FK) · text · status · created_at\n"
         "оценки: rhyme_score · structure_score · style_score · individuality_score · vibe_score · total_score",
         accent=ACCENT_RED)


# ---------- Слайд 6: Ключевые алгоритмы (стиль v2, диапазон 6-90) ----------


def slide_algorithm(prs, total):
    s = blank(prs)
    chrome(s, "Ключевые алгоритмы: расчёт оценки", 6, total)

    # Формула
    fx, fy = Inches(0.55), Inches(1.45)
    fw = Inches(12.2)
    add_rounded(s, fx, fy, fw, Inches(1.85), fill=NAVY, line=None, radius=0.06)
    add_text(s, fx, fy + Inches(0.12), fw, Inches(0.4),
             "Итоговая оценка рецензии", font=BODY_FONT, size=15, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, fx, fy + Inches(0.5), fw, Inches(0.9),
             "Score = (R + S + I + Ch) × 1.4 × A", font=SERIF_FONT, size=46,
             bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, fx, fy + Inches(1.42), fw, Inches(0.35),
             "округляется до целого; фактический диапазон итоговой оценки: 6 – 90",
             font=BODY_FONT, size=13, italic=True, color=GRAY_50, align=PP_ALIGN.CENTER)

    # Карточки критериев
    crit = [
        ("R", "Рифмы / образы", "1 – 10", NAVY),
        ("S", "Структура / ритмика", "1 – 10", NAVY),
        ("I", "Реализация стиля", "1 – 10", NAVY),
        ("Ch", "Индивидуальность", "1 – 10", NAVY),
        ("A", "Атмосфера / вайб (×)", "1.0 – 1.6", ACCENT_RED),
    ]
    cw = Inches(2.32)
    gap = Inches(0.16)
    start_x = (SLIDE_W - cw * 5 - gap * 4) // 2
    cy = Inches(3.65)
    ch = Inches(1.95)
    for i, (sym, name, rng, ccol) in enumerate(crit):
        x = start_x + (cw + gap) * i
        add_rounded(s, x, cy, cw, ch, fill=WHITE, line=GRAY_BORDER)
        add_circle(s, x + (cw - Inches(0.95)) // 2, cy + Inches(0.22), Inches(0.95), fill=ccol)
        add_text(s, x + (cw - Inches(0.95)) // 2, cy + Inches(0.22), Inches(0.95), Inches(0.95),
                 sym, font=TITLE_FONT, size=26, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.1), cy + Inches(1.25), cw - Inches(0.2), Inches(0.4),
                 name, font=BODY_FONT, size=14, bold=True, color=NAVY_DARK,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.1), cy + Inches(1.62), cw - Inches(0.2), Inches(0.3),
                 rng, font=BODY_FONT, size=13, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # Нижняя плашка-сравнение
    by = Inches(5.85)
    add_rounded(s, Inches(0.55), by, Inches(12.2), Inches(1.05), fill=GRAY_50, line=GRAY_BORDER)
    add_text(s, Inches(0.85), by + Inches(0.13), Inches(5.6), Inches(0.35),
             "Аналоги (RYM, AllMusic, VK):", font=BODY_FONT, size=14, bold=True, color=NAVY_DARK)
    add_text(s, Inches(0.85), by + Inches(0.5), Inches(5.6), Inches(0.45),
             "одна звёздная шкала 1–5 → высокая субъективность, низкая информативность.",
             font=BODY_FONT, size=13, color=TEXT_MUTED, line_spacing=1.05)
    add_text(s, Inches(6.9), by + Inches(0.13), Inches(5.6), Inches(0.35),
             "Разрабатываемая система:", font=BODY_FONT, size=14, bold=True, color=ACCENT_RED)
    add_text(s, Inches(6.9), by + Inches(0.5), Inches(5.6), Inches(0.45),
             "5 независимых критериев → структурированная оценка и наглядное сравнение релизов.",
             font=BODY_FONT, size=13, color=TEXT_MUTED, line_spacing=1.05)


# ---------- Слайд 7: Архитектура и DevOps (нативно) ----------


def flow_box(slide, x, y, w, h, title, subtitle, fill=WHITE, tcol=NAVY_DARK):
    add_rounded(slide, x, y, w, h, fill=fill, line=GRAY_BORDER if fill == WHITE else None)
    add_text(slide, x + Inches(0.15), y + Inches(0.12), w - Inches(0.3), Inches(0.4),
             title, font=BODY_FONT, size=14, bold=True, color=tcol,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, x + Inches(0.12), y + Inches(0.48), w - Inches(0.24), h - Inches(0.55),
                 subtitle, font=BODY_FONT, size=12,
                 color=TEXT_MUTED if fill == WHITE else GRAY_50,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.05)


def arrow_down(slide, cx, y, h, color=NAVY_LIGHT):
    a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, cx - Inches(0.13), y, Inches(0.26), h)
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background()
    a.shadow.inherit = False


def slide_devops(prs, total):
    s = blank(prs)
    chrome(s, "Архитектура системы и DevOps", 7, total)

    # ===== Левая колонка: архитектура =====
    lx = Inches(0.55)
    lw = Inches(5.75)
    add_text(s, lx, Inches(1.25), lw, Inches(0.4), "Архитектура приложения",
             font=TITLE_FONT, size=17, bold=True, color=NAVY_DARK)
    bx = lx + Inches(0.95)
    bw = Inches(3.85)
    bh = Inches(0.78)
    ys = [Inches(1.8), Inches(2.92), Inches(4.04), Inches(5.16)]
    boxes = [
        ("Браузер · React SPA", "интерфейс, axios → REST", WHITE, NAVY_DARK),
        ("Nginx (reverse-proxy)", "статика + проксирование /api", NAVY, WHITE),
        ("Backend · Go + Gin", "REST API, бизнес-логика, GORM", WHITE, NAVY_DARK),
        ("PostgreSQL", "каталог, рецензии, лайки, связи", NAVY, WHITE),
    ]
    for i, (t, sub, f, tc) in enumerate(boxes):
        flow_box(s, bx, ys[i], bw, bh, t, sub, fill=f, tcol=tc)
        if i < len(boxes) - 1:
            arrow_down(s, bx + bw // 2, ys[i] + bh + Inches(0.02), Inches(0.3))

    # ===== Правая колонка: CI/CD =====
    rx = Inches(6.95)
    rw = Inches(5.8)
    add_text(s, rx, Inches(1.25), rw, Inches(0.4), "CI/CD · GitHub Actions",
             font=TITLE_FONT, size=17, bold=True, color=NAVY_DARK)
    stages = [
        ("1 · Backend", "go vet · go test · go build"),
        ("2 · Frontend", "npm install · npm run build"),
        ("3 · Compose smoke", "поднять prod-стек, проверить /healthz"),
        ("4 · Docker / GHCR", "сборка и публикация образов"),
    ]
    sx = rx + Inches(0.35)
    sw = Inches(5.1)
    sh = Inches(0.78)
    sy = [Inches(1.8), Inches(2.92), Inches(4.04), Inches(5.16)]
    for i, (t, sub) in enumerate(stages):
        add_rounded(s, sx, sy[i], sw, sh, fill=GRAY_50, line=GRAY_BORDER)
        add_circle(s, sx + Inches(0.16), sy[i] + Inches(0.19), Inches(0.4), fill=NAVY)
        add_text(s, sx + Inches(0.16), sy[i] + Inches(0.19), Inches(0.4), Inches(0.4),
                 str(i + 1), font=TITLE_FONT, size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, sx + Inches(0.7), sy[i] + Inches(0.1), sw - Inches(0.85), Inches(0.32),
                 t.split(" · ")[1] if " · " in t else t, font=BODY_FONT, size=14,
                 bold=True, color=NAVY_DARK)
        add_text(s, sx + Inches(0.7), sy[i] + Inches(0.42), sw - Inches(0.85), Inches(0.3),
                 sub, font=BODY_FONT, size=12.5, color=TEXT_MUTED)
        if i < len(stages) - 1:
            arrow_down(s, sx + sw // 2, sy[i] + sh + Inches(0.02), Inches(0.3))

    # ===== чипы технологий внизу (центрированы) =====
    chips = ["Docker Compose", "Multi-stage build", "Healthcheck", "GHCR registry", "Graceful shutdown"]
    cy = Inches(6.25)
    chip_gap = Inches(0.2)
    widths = [Inches(0.35 + 0.105 * len(c)) for c in chips]
    total = sum(widths, Emu(0)) + chip_gap * (len(chips) - 1)
    cx = (SLIDE_W - total) // 2
    for c, cw in zip(chips, widths):
        add_rounded(s, cx, cy, cw, Inches(0.5), fill=NAVY, line=None, radius=0.3)
        add_text(s, cx, cy, cw, Inches(0.5), c, font=BODY_FONT, size=12, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += cw + chip_gap


# ---------- Слайд 8: Полученные результаты (стиль v1) ----------


def slide_results(prs, total):
    s = blank(prs)
    chrome(s, "Полученные результаты", 8, total)

    stats = [("30+", "эндпоинтов REST API"), ("5", "таблиц базы данных"),
             ("4", "стадии CI-пайплайна"), ("1", "команда запуска")]
    sw = Inches(2.95)
    gap = Inches(0.2)
    start_x = (SLIDE_W - sw * 4 - gap * 3) // 2
    top = Inches(1.3)
    sh = Inches(1.45)
    for i, (kpi, label) in enumerate(stats):
        x = start_x + (sw + gap) * i
        add_rounded(s, x, top, sw, sh, fill=NAVY, line=None)
        add_text(s, x, top + Inches(0.1), sw, Inches(0.85), kpi, font=TITLE_FONT,
                 size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, top + Inches(0.98), sw, Inches(0.4), label, font=BODY_FONT,
                 size=12, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    list_top = Inches(3.05)
    # Левая колонка — функциональность
    add_text(s, Inches(0.55), list_top, Inches(6.0), Inches(0.4),
             "Функциональность", font=TITLE_FONT, size=19, bold=True, color=NAVY_DARK)
    left_items = [
        "Авторизация по подписанным сессионным токенам, роли admin / user",
        "Многокритериальная оценка релизов и оптимистичные лайки",
        "Лента рецензий с фильтром «Все / Подписки»",
        "Конвейер модерации: pending → approved / rejected",
        "Профиль с уровнями, бейджами и социальным графом",
        "Музыкальный каталог: альбомы, треки, жанры",
    ]
    # Правая колонка — инженерная часть
    add_text(s, Inches(6.95), list_top, Inches(6.0), Inches(0.4),
             "Инфраструктура и качество", font=TITLE_FONT, size=19, bold=True, color=NAVY_DARK)
    right_items = [
        "Контейнеризация: Docker Compose (dev / prod / deploy)",
        "CI: go vet · go test · go build + frontend build",
        "Smoke-тест prod-стека и публикация образов в GHCR",
        "Multi-stage сборка, healthcheck, graceful shutdown",
        "Идемпотентный сидер демо-данных",
        "Запуск одной командой, готовность к деплою на VPS",
    ]
    for col_x, items in ((Inches(0.55), left_items), (Inches(6.95), right_items)):
        y = list_top + Inches(0.55)
        for txt in items:
            add_text(s, col_x + Inches(0.1), y, Inches(0.35), Inches(0.4), "▸",
                     font=BODY_FONT, size=14, bold=True, color=ACCENT_RED)
            add_text(s, col_x + Inches(0.45), y, Inches(5.5), Inches(0.55), txt,
                     font=BODY_FONT, size=15, color=TEXT_DARK, line_spacing=1.05)
            y += Inches(0.6)


# ---------- Слайд 9: Спасибо ----------


def slide_thanks(prs):
    s = blank(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY_DARK)
    add_rect(s, 0, Inches(3.5), SLIDE_W, Inches(0.06), fill=ACCENT_RED)
    logo(s, Inches(6.16), Inches(0.7), Inches(2.1), badge=True, pad=Inches(0.16))
    add_text(s, Inches(0.5), Inches(3.95), Inches(12.3), Inches(1.4),
             "СПАСИБО ЗА ВНИМАНИЕ!", font=TITLE_FONT, size=60, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
             f"{AUTHOR_FULL}  ·  группа {GROUP}  ·  {CITY_YEAR}",
             font=BODY_FONT, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    total = 9
    slide_title(prs)
    slide_goal(prs, total)
    slide_tasks(prs, total)
    slide_analytics(prs, total)
    slide_db(prs, total)
    slide_algorithm(prs, total)
    slide_devops(prs, total)
    slide_results(prs, total)
    slide_thanks(prs)
    prs.save(str(OUTPUT))
    print(f"saved: {OUTPUT}")


if __name__ == "__main__":
    main()
