# -*- coding: utf-8 -*-
"""Добавляет инфографику (карточки сущностей) на слайд 5 'Логическая модель БД'.
Открывает presentation_v3.pptx, НЕ меняет другие слайды, сохраняет как presentation_v4.pptx."""
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = "presentation_v3.pptx"
DST = "presentation_v4.pptx"

# --- палитра деки ---
NAVY        = RGBColor(0x1F, 0x3A, 0x8A)
NAVY_DARK   = RGBColor(0x15, 0x29, 0x66)
ACCENT_RED  = RGBColor(0xC8, 0x10, 0x2E)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_BORDER = RGBColor(0xD8, 0xDE, 0xEB)
TEXT_MUTED  = RGBColor(0x5A, 0x6B, 0x8C)

shutil.copyfile(SRC, DST)
prs = Presentation(DST)
slide = prs.slides[4]   # слайд 5


def add_text(x, y, w, h, text, size, color, bold=False, font="Calibri",
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.name = font; r.font.color.rgb = color
    return tb


def add_round(x, y, w, h, fill, line=None, radius=0.10, line_w=0.75, shadow=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    if not shadow:
        shp.shadow.inherit = False
    return shp


# --- геометрия колонки ---
x0 = 7.50
col_w = 5.45
card_h = 0.74
gap = 0.14
y_start = 2.36

# мини-заголовок колонки
add_text(x0 + 0.06, 1.90, col_w, 0.32, "СОСТАВ МОДЕЛИ",
         size=12, color=TEXT_MUTED, bold=True)

entities = [
    ("users",   "Учётные записи, роли и профили пользователей", NAVY),
    ("albums",  "Каталог релизов: метаданные и обложки",        NAVY),
    ("tracks",  "Композиции альбома и их атрибуты",              NAVY),
    ("genres",  "Справочник музыкальных жанров (M:N)",           NAVY),
    ("reviews", "Оценка по 5 критериям + статус модерации",      ACCENT_RED),
]

tag_w = 1.45
tag_h = 0.40
for i, (name, desc, accent) in enumerate(entities):
    y = y_start + i * (card_h + gap)
    # карточка
    add_round(x0, y, col_w, card_h, fill=WHITE, line=GRAY_BORDER,
              radius=0.12, line_w=0.75, shadow=True)
    # синий/красный тег с именем таблицы (как шапки таблиц на диаграмме)
    tag = add_round(x0 + 0.20, y + (card_h - tag_h) / 2, tag_w, tag_h,
                    fill=accent, line=None, radius=0.22, shadow=False)
    tf = tag.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = name
    r.font.size = Pt(13); r.font.bold = True
    r.font.name = "Consolas"; r.font.color.rgb = WHITE
    # описание
    dx = x0 + 0.20 + tag_w + 0.22
    add_text(dx, y + 0.04, x0 + col_w - dx - 0.20, card_h - 0.08, desc,
             size=11, color=TEXT_MUTED, anchor=MSO_ANCHOR.MIDDLE)

prs.save(DST)
print("saved", DST)
print("last card bottom:", round(y_start + 4 * (card_h + gap) + card_h, 3))
